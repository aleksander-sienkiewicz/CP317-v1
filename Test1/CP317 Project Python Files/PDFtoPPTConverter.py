import fitz
import pdfplumber
from pypdf import PdfReader
from pptx.util import Inches
from pptx import Presentation

class PDFtoPPTConverter:
    def __init__(self, name: str, encoding: str ="utf-8", reverse: bool =False):
        self.name = name
        self.file = PdfReader(name)
        self.len = len(self.file.pages)
        self.encoding = encoding
        self.reverse = reverse
        self.pages = []
        self.tables = {x: [] for x in range(self.len)}
        

    def extractTextAndImages(self):
        pages = []
        for x in range(self.len):
            # text stuff
            text = self.file.pages[x].extract_text().encode(self.encoding, errors = "ignore").decode(self.encoding, errors = "ignore")
            pages.append([text, 0])
            #image stuff
            with fitz.open(self.name) as temp_open:
                images = temp_open[x].get_images()
                for index, y in enumerate(images, start=1):
                    xref = y[0]
                    pix = fitz.Pixmap(self.file, xref)
                    pix.save("images/page_%s-image_%s.png" % (x, index))
                    pages[x][1] += 1
        
        if(self.reverse):
            pages.reverse()
        self.pages = pages
    def extractImages(self):
            
            for x in range(self.len):
            #image stuff
            
                images = self.file[x].get_images()
                for index, y in enumerate(images, start=1):
                    xref = y[0]
                    pix = fitz.Pixmap(self.file, xref)
                    pix.save("images/page_%s-image_%s.png" % (x, index))

    def extractTables(self):
        
        i = 0
        with pdfplumber.open(self.name) as pdf_file:
            for page in pdf_file.pages:
                
                self.tables[i] = page.extract_tables()
                print(self.tables[i])
                i += 1
        

    def create_ppt(self, name: str):
        prs = Presentation()
        i = 0
        for page in self.pages:
            second_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(second_layout)
            title = slide.shapes.title
            title.text = "Page " + str(i + 1)
            slide.placeholders[1].text = page[0]
            # adding tables first
            j = 1
            for table in self.tables[i]:
                print(table)
                sectionH_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(sectionH_layout)
                slide.shapes[0].text = "Table " + str(j)
                shape = slide.shapes.add_table(len(table), len(table[0]), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5))
                ppt_table = shape.table
                for row in range(len(table)):
                    for column in range(len(table[0])):
                        curr_cell = ppt_table.cell(row, column)
                        curr_cell.text = table[row][column]
                j += 1
            # adding images next
            for i2 in range(page[1]):
                image_location = f"/images/page_{i}-image_{i2}.png"
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)
                picture = slide.shapes.add_picture(image_location, Inches(1), Inches(1), Inches(7), Inches(7))
            i += 1
        prs.save(name)
            
                        
                
                
            
            
            
            


if __name__ == "__main__":
    i = PDFtoPPTConverter("test4.pdf")
    i.extractTables()
    i.extractTextAndImages()
    i.create_ppt("test.pptx")
    
    