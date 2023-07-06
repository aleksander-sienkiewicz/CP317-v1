import fitz
import tabula
import pandas
from sympy import preview
from pptx import Presentation



class articleDissasembler:
    def __init__(self, fileName):
         self.fileName = fileName
         self.file = fitz.open(fileName)
         self.length = len(self.file)

    def extractTextAndImages(self):
        pages = []
        for x in range(0, self.length):
            # text stuff
            text = self.file[x].get_text().encode("utf8")
            decoded = text.decode()
            pages.append(decoded)
            #image stuff
            images = self.file[x].get_images()
            for index, y in enumerate(images, start=1):
                xref = y[0]
                pix = fitz.Pixmap(self.file, xref)
                pix.save("images/page_%s-image_%s.png" % (x, index))
        return(pages)
    def extractTables(self):
        #get images and text in each page
        page = []
        tablepage = []

        #gets tables in each page
        for i in (range(1,(self.length+1))):
            stri = str(i)
            table = tabula.read_pdf(self.fileName, pages=stri, multiple_tables=True, stream='true')

            for i in table:
                page.append(i.to_records())
            tablepage.append(page)
            page = []

        return(tablepage)
    def getEquations(self):
        a = self.extractTextAndImages()
        temp = ""
        counter = 0
        line = 0
        for i in range(0,len(a[4])-1):
            if (a[4][i] != '\n'):
                counter +=1
            else:
                print(counter, line)
                line += 1
                counter = 0
        # page = self.file.new_page()
        # text = '\xe2\x8e\xaa\xe2\x8e\xa9\n\xe2\x8e\xaa\xe2\x8e\xa8\n\xe2\x8e\xa7\n\xe2\x89\xa0\n\xe2\x88\x92\n=\n= \xe2\x88\x91\n\xe2\x89\xa0\nj\n'
        # p = fitz.Point(50, 72) 

        # rc = page.insert_text(p,  # bottom-left of 1st char
        #                     text,  # the text (honors '\n')
        #                     fontname = "helv",  # the default font
        #                     fontsize = 11,  # the default font size
        #                     rotate = 0,  # also available: 90, 180, 270
        #                     )
        # print("%i lines printed on page %i." % (rc, page.number))

        # self.file.save("output2.pdf")

         

# a = articleDissasembler("tablesample.pdf")
# b = a.extractTables()

# print(b)

preview(r'$$\int_0^1 e^x\,dx$$', viewer='file', filename='test.png', euler=False)


    




# prs = Presentation()
# bullet_slide_layout = prs.slide_layouts[1]

# slide = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes

# title_shape = shapes.title
# body_shape = shapes.placeholders[1]

# title_shape.text = 'Adding a Bullet Slide'
# tf = body_shape.text_frame
# tf.text = t

# prs.save('./test.pptx')
