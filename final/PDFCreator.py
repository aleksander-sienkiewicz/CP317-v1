import os
from transformers import pipeline
from pptx import Presentation
from pptx.util import Inches, Pt

class createPowerPoint():
    def __init__(self, fileName):
         self.output_file = fileName
         self.prs = Presentation()

    def add_slides(self,summaries, images, tables, page_index):
        # Create a PowerPoint presentation
        # Set the width and height of a slide
        currpage = 0 # is the current page in corrolation to the summaires
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        self.prs.slide_width = slide_width
        self.prs.slide_height = slide_height
        for i, summary in enumerate(summaries):
            self.add_summary(summary, i, currpage) # adds a summary slide
            while (currpage < len(page_index) and i == page_index[currpage]): # checks if the slide corresponds to the end of the page. If it does add images and tables
                 for x in range(0, len(images[currpage])): # add slides for all images on the current page
                      self.add_image(images[currpage][x], currpage,x)
                 for x in range(0, len(tables[currpage])): # add tables for all tables on the current page
                      self.add_table(tables[currpage][x], currpage, x)
                 currpage += 1
        # Save the presentation
        self.prs.save(self.output_file)

    def add_summary(self, summary, i, currpage):
        # Add a slide with a title and content
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            # Set the title and content
            title.text = f"Page: {currpage+1}"
            content.text = summary

    def add_image(self, imgpath, pagenum, imagenum):
        # Add a slide with a title and content
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            # Set the title and content
            left  = int(self.prs.slide_width * 0.15)
            top   = int(self.prs.slide_height * 0.17)
            width = int(self.prs.slide_width * 0.7)
            title.text = f"Page: {pagenum+1} Image: {imagenum+1}"
            slide.shapes.add_picture(imgpath, left, top, width)
            os.remove(imgpath)
        
    def add_table(self, table, pagenum, tablenum):
         slide_layout = self.prs.slide_layouts[1]
         slide = self.prs.slides.add_slide(slide_layout)
         title = slide.shapes.title
         # Set the title and content
         title.text = f"Page: {pagenum+1} Table: {tablenum+1}"
         x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5.5)
         if (len(table) != 0):
          shape = slide.shapes.add_table(len(table), len(table[0])-1, x, y, cx, cy)
          tab = shape.table
          # add table headers
          headers = table.dtype.names
          for i in range(1, len(headers)):
               cell = tab.cell(0,i-1)
               cell.text = headers[i]
          #add table contents
          for rows in range(0,len(table)):
               for cols in range(1,len(table[rows])):
                    cell = tab.cell(rows,cols-1)
                    cell.text = str(table[rows][cols])



    def save_ppt(self):
        self.prs.save(self.output_file)