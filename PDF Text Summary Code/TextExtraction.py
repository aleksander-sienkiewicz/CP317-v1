import os
import fitz
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter, PDFPageAggregator
from pdfminer.layout import LAParams, LTTextBox, LTTextLine, LTTextBoxHorizontal, LTImage, LTChar, LTAnno, LTText
from pdfminer.pdfpage import PDFPage
from io import StringIO
from pptx import Presentation
from pptx.util import Inches
def extract_content_from_pdf(pdf_path):
    # Create a PDF resource manager object
    resource_manager = PDFResourceManager()

    # Create a PDF interpreter object
    page_interpreter = PDFPageInterpreter(resource_manager, PDFPageAggregator(resource_manager, laparams=LAParams()))

    # Initialize empty lists to store the extracted content
    text_content = []
    image_content = []
    equation_content = []
    table_content = []

    # Open the PDF file
    with open(pdf_path, 'rb') as pdf_file:
        # Iterate over PDF pages
        for page in PDFPage.get_pages(pdf_file):
            # Process the page
            page_interpreter.process_page(page)

            # Get the page layout
            layout = page_interpreter.device.get_result()

            # Iterate over layout objects
            for obj in layout:
                # Check if the object is a text box or text line
                if isinstance(obj, (LTTextBox, LTTextLine, LTTextBoxHorizontal)):
                    # Extract the text from the object
                    text = obj.get_text().strip()
                    if text:
                        # Append the text to the result
                        text_content.append(text)

                # Check if the object is an image
                elif isinstance(obj, LTImage):
                    # Append the image to the result
                    image_content.append(obj)

                # Check if the object is an equation
                elif isinstance(obj, (LTChar, LTAnno)) and obj.get_text().startswith('$'):
                    # Extract the equation from the object
                    equation = obj.get_text().strip()
                    if equation:
                        # Append the equation to the result
                        equation_content.append(equation)

                # Check if the object is a table
                elif isinstance(obj, LTText) and '|' in obj.get_text():
                    # Extract the table from the object
                    table = []
                    for line in obj.get_text().split('\n'):
                        row = [cell.strip() for cell in line.split('|')]
                        if row:
                            table.append(row)
                    if table:
                        # Append the table to the result
                        table_content.append(table)

    # Return the extracted content
    return text_content, image_content, equation_content, table_content
def extract_text(self):
        file = fitz.open(self.file_path)
        final = ""
        endOfPage = []
        for x in range(0, len(file)):
            text = file[x].get_text().encode("utf8")
            decoded = text.decode()
            final += " "
            final += decoded
            endOfPage.append(len(final)// 1024)

        return(final, endOfPage)
def create_powerpoint_from_pdf(pdf_path, pptx_path):
    # Extract the content from the PDF file
    text_content, image_content, equation_content, table_content = extract_content_from_pdf(pdf_path)

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Add a slide for each page of the PDF file
    for i in range(len(text_content)):
        # Add a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Add the text content to the slide
        text_frame = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        text_frame.text = text_content[i]

        # Add the image content to the slide
        for image in image_content:
            image_file = 'image.png'
            with open(image_file, 'wb') as f:
                f.write(image.stream.get_rawdata())
            slide.shapes.add_picture(image_file, Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            os.remove(image_file)

        # Add the equation content to the slide
        for equation in equation_content:
            text_frame = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            text_frame.text = equation

        # Add the table content to the slide
        for table in table_content:
            rows = len(table)
            cols = len(table[0])
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(6))
            for i in range(rows):
                for j in range(cols):
                    table_shape.table.cell(i, j).text = table[i][j]

    # Save the PowerPoint presentation
    prs.save(pptx_path)
    

pdf_path = 'uploads/test.pdf'
pptx_path = 'uploads/output.pptx'
create_powerpoint_from_pdf(pdf_path, pptx_path)