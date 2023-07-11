import PyPDF2
import fitz
import pdfplumber
from transformers import pipeline
import nltk
from pptx import Presentation
from pptx.util import Inches, Pt

class createPowerPoint():
    def __init__(self, fileName):
         self.file_path = fileName
         self.prs = Presentation()

    def extract_text(self):
        file = fitz.open(self.file_path)
        final = ""
        endOfPage = []
        for x in range(0, len(file)):
            text = file[x].get_text().encode("utf8")
            decoded = text.decode()
            final += " "
            final += decoded
            endOfPage.append(len(final)// 1024)

        return(final, endOfPage)

    def summarize_pdf(self, text):
       # Read the PDF document
            
        # Tokenize the text into sentences
        sentences = nltk.sent_tokenize(text)
        # Combine the sentences into a single string
        document = ' '.join(sentences)
        # Initialize the BART summarization pipeline
        summarizer = pipeline('summarization')
        # Calculate how many chunks we need
        chunks = len(document) // 1024 + (len(document) % 1024 > 0)
        # Create a list for summaries
        summaries = []
        # Summarize each chunk and add to the list of summaries
        for i in range(chunks):
            chunk = document[i*1024:(i+1)*1024]
            summary = summarizer(chunk, max_length=60, min_length=30, do_sample=False)
            # Extract the summarized text
            summarized_text = summary[0]['summary_text']
            summaries.append(summarized_text)
            
        return summaries



    def add_summaries(self,summaries, output_file):
        # Create a PowerPoint presentation
        # Set the width and height of a slide
        slide_width = Inches(10)
        slide_height = Inches(7.5)
        self.prs.slide_width = slide_width
        self.prs.slide_height = slide_height
        for i, summary in enumerate(summaries):
            # Add a slide with a title and content
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            # Set the title and content
            title.text = f"Slide {i+1}"
            content.text = summary

        # Save the presentation
        self.prs.save(output_file)


# Usage example

a = createPowerPoint("test2.pdf")
text, endOfPage = a.extract_text()
print(endOfPage)
# summaries = a.summarize_pdf(text)

# a.add_summaries(summaries, "./text.pptx")
