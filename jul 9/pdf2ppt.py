import os
import PyPDF2
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import PDFPageAggregator
from pdfminer.layout import LAParams, LTImage, LTTextBox, LTTextLine, LTChar, LTAnno, LTFigure, LTTextBoxHorizontal
from pptx import Presentation
from pptx.util import Inches

class ExtractContent():
    def __init__(self, fileName):
         self.fileName = fileName
         self.text_content = []
         self.image_content = []
         self.equation_content = []
         self.table_content = []

    def extract_content_from_pdf(self):
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfFileReader(open(self.pdf_path, 'rb'))

        # Create a PDF resource manager object
        resource_manager = PDFResourceManager()

        # Create a PDF interpreter object
        page_interpreter = PDFPageInterpreter(resource_manager, PDFPageAggregator(resource_manager, laparams=LAParams()))

        # Iterate over PDF pages
        for page_num in range(pdf_reader.getNumPages()):
            # Get the page object
            page = pdf_reader.getPage(page_num)

            # Process the page
            page_interpreter.process_page(page)

            # Get the page layout
            layout = page_interpreter.device.get_result()

            # Iterate over layout objects
            for obj in layout:
                # Check if the object is a text box or text line
                if isinstance(obj, (LTTextBox, LTTextLine, LTTextBoxHorizontal)):
                    # Extract the text from the object
                    text = obj.get_text().strip()
                    if text:
                        # Append the text to the result
                        self.text_content.append(text)

                # Check if the object is an image
                elif isinstance(obj, LTImage):
                    # Append the image to the result
                    self.image_content.append(obj)

                # Check if the object is an equation
                elif isinstance(obj, (LTChar, LTAnno)) and obj.get_text().startswith('$'):
                    # Extract the equation from the object
                    equation = obj.get_text().strip()
                    if equation:
                        # Append the equation to the result
                        self.equation_content.append(equation)

                # Check if the object is a table
                # elif isinstance(obj, LTTable):
                #     # Extract the table from the object
                #     table = []
                #     for row in obj:
                #         table.append([cell.get_text().strip() for cell in row])
                #     if table:
                #         # Append the table to the result
                #         table_content.append(table)

        # Return the extracted content

    # def create_powerpoint_from_pdf(pdf_path, pptx_path):
    #     # Extract the content from the PDF file
    #     text_content, image_content, equation_content, table_content = extract_content_from_pdf(pdf_path)

    #     # Create a new PowerPoint presentation
    #     prs = Presentation()

    #     # Add a slide for each page of the PDF file
    #     for i in range(len(text_content)):
    #         # Add a new slide
    #         slide = prs.slides.add_slide(prs.slide_layouts[1])

    #         # Add the text content to the slide
    #         text_frame = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    #         text_frame.text = text_content[i]

    #         # Add the image content to the slide
    #         for image in image_content:
    #             image_file = 'image.png'
    #             with open(image_file, 'wb') as f:
    #                 f.write(image.stream.get_rawdata())
    #             slide.shapes.add_picture(image_file, Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    #             os.remove(image_file)

    #         # Add the equation content to the slide
    #         for equation in equation_content:
    #             text_frame = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    #             text_frame.text = equation

    #         # Add the table content to the slide
    #         # for table in table_content:
    #         #     rows = len(table)
    #         #     cols = len(table[0])
    #         #     table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(9), Inches(6))
    #         #     for i in range(rows):
    #         #         for j in range(cols):
    #         #             table_shape.table.cell(i, j).text = table[i][j]

    #     # Save the PowerPoint presentation
    #     prs.save(pptx_path)

# Example usage
pdf_path = 'test.pdf'
output_ppt = 'output.pptx'
a = ExtractContent("pdf_path")
a.extract_content_from_pdf()
